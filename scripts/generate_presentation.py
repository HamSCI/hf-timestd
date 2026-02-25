#!/usr/bin/env python3
"""
Generate HamSCI 2026 PowerPoint presentation from template.

Uses the HamSCI oral presentation template and populates it with
our 11-slide outline, figures, and speaker notes.

Usage:
    python3 scripts/generate_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

TEMPLATE = os.path.join(os.path.dirname(__file__), '..', 'docs',
                        'HamSCI Oral PresentationTemplate.pptx')
FIGURES = os.path.join(os.path.dirname(__file__), '..', 'docs', 'figures')
OUTPUT = os.path.join(os.path.dirname(__file__), '..', 'docs',
                      'HamSCI_2026_AC0G.pptx')

# Slide dimensions (from template): 13.333 x 7.5 inches
SLIDE_W = 13.333
SLIDE_H = 7.5


def fig(name):
    """Return full path to a figure."""
    return os.path.join(FIGURES, name)


def set_notes(slide, text):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if color:
        p.font.color.rgb = RGBColor(*color)
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=None, font_name='Calibri'):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.space_after = Pt(4)
        p.level = 0
        if color:
            p.font.color.rgb = RGBColor(*color)
    return tf


def add_table(slide, left, top, width, height, rows, col_widths=None):
    """Add a table to a slide. rows = list of lists of strings."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(13)
    return table


def build_presentation():
    prs = Presentation(TEMPLATE)

    # Remove all template slides (we'll build fresh using layouts)
    ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    for sldId in list(prs.slides._sldIdLst):
        rId = sldId.get(ns + 'id')
        prs.slides._sldIdLst.remove(sldId)
        prs.part.drop_rel(rId)

    title_layout = prs.slide_layouts[0]      # Title Slide
    content_layout = prs.slide_layouts[1]     # Title and Content
    subtitle_layout = prs.slide_layouts[2]    # Subtitle

    # =================================================================
    # SLIDE 1: Title
    # =================================================================
    slide = prs.slides.add_slide(title_layout)
    slide.placeholders[0].text = (
        "With an RX888 and a GPSDO,\n"
        "What Kind of Ionospheric Science Can We Do?")
    slide.placeholders[1].text = (
        "Michael James Hauan (AC0G)\n"
        "EM38, Central Missouri\n\n"
        "4 stations · 9 frequencies · 17 simultaneous paths\n"
        "Single RX888 SDR + Leo Bodnar GPSDO via KA9Q-radio")
    set_notes(slide,
        "This talk is organized around one question: if you have an RX888 "
        "and a GPSDO listening to time standard stations, what ionospheric "
        "science can you actually extract? The answer depends on your timing "
        "infrastructure, and it turns out the science payoff is much larger "
        "than you might expect.")

    # =================================================================
    # SLIDE 2: Phenomena Ladder (fig15)
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "The Phenomena Ladder — What Each Tier Unlocks"

    # Full-width figure
    slide.shapes.add_picture(
        fig('fig15_phenomena_ladder.png'),
        Inches(0.3), Inches(0.7), Inches(12.7), Inches(6.5))

    set_notes(slide,
        "Here's the organizing framework. On the left is what hardware you "
        "have. On the right is what ionospheric phenomena you can observe. "
        "The GPSDO — around 160 dollars for a Leo Bodnar — is the single "
        "most important upgrade. It gives you a 1 part-per-billion sample "
        "clock, which makes carrier-phase measurements metrologically "
        "coherent. That unlocks Doppler, dTEC/dt, and scintillation — all "
        "rate measurements that don't need to know what time it is, only "
        "that consecutive samples are exactly spaced. The absolute delay "
        "products — D_clock, mode identification, group-delay TEC — need to "
        "know when in UTC each tick arrived. That's either an external PPS "
        "reference or, as I'll show, something we can recover from the time "
        "signals themselves.")

    # =================================================================
    # ACT 2 divider (not a separate slide — just a visual cue)
    # SLIDE 3: TickEdgeDetector
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "TickEdgeDetector — The Measurement Engine"

    # Left column: description
    add_bullet_list(slide, 0.4, 0.9, 6.0, 5.5, [
        "50–57 ticks per minute, per channel",
        "Matched-filter correlation with station-specific templates",
        "Front-edge back-calculation + sub-sample interpolation",
        "SNR-weighted robust median ensemble",
        "",
        "From each tick:",
        "  • Timing error (AM domain)",
        "  • Carrier phase (IQ domain)",
        "  • SNR",
        "",
        "From the minute ensemble:",
        "  • D_clock — needs UTC (Tier 3+)",
        "  • Doppler — needs only GPSDO (Tier 2)",
        "  • Mean SNR",
    ], font_size=15)

    # Right column: evidence table
    add_table(slide, 6.8, 1.0, 6.0, 2.5, [
        ['Channel', 'Edges/min', 'SNR (dB)', 'D_clock (ms)', 'Doppler'],
        ['CHU 7.85', '51', '27.3', '+1.96', '99.8%'],
        ['SHARED 10.0', '57', '7.9', '−0.02', '99.7%'],
        ['WWV 20.0', '57', '7.9', '−0.05', '99.9%'],
    ])

    add_textbox(slide, 6.8, 3.8, 6.0, 1.5,
                "~24K tick timing records/day\n"
                "~850K per-tick phase records/day\n"
                "9 channels, 24/7",
                font_size=14, color=(0x44, 0x44, 0x44))

    # Detection summary figure
    if os.path.exists(fig('fig8_detection_summary.png')):
        slide.shapes.add_picture(
            fig('fig8_detection_summary.png'),
            Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.5))

    set_notes(slide,
        "The measurement engine is a tick edge detector inspired by the "
        "ntpd WWV refclock driver. For every second of the minute, it "
        "cross-correlates a station-specific template against the IQ data, "
        "finds the front edge with sub-sample precision, and extracts both "
        "the timing error and the carrier phase. From 50 to 57 ticks, we "
        "build a robust median ensemble. The carrier phase across the minute "
        "gives Doppler — and notice, that only needs a stable sample clock. "
        "It doesn't need to know what time it is. But D_clock — the absolute "
        "propagation delay — does need UTC. That's the bridge we build next.")

    # =================================================================
    # SLIDE 4: UTC Recovery
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "UTC Recovery — Dual Kalman Fusion"

    # Left: metrological ladder figure
    if os.path.exists(fig('fig1_metrological_ladder.png')):
        slide.shapes.add_picture(
            fig('fig1_metrological_ladder.png'),
            Inches(0.3), Inches(0.9), Inches(6.2), Inches(4.5))

    # Right: key results
    add_textbox(slide, 6.8, 0.9, 6.0, 0.5,
                "The time signals tell us what time it is.",
                font_size=18, bold=True, color=(0x15, 0x65, 0xc0))

    add_table(slide, 6.8, 1.6, 5.8, 2.2, [
        ['Source', 'Offset vs GPS', 'Bound (±)'],
        ['Internet NTP', '−6.0 ms', '±19 ms'],
        ['HF TSL1 (geometric)', '+1.1 ms', '±2.0 ms'],
        ['HF TSL2 (ionospheric)', '+0.8 ms', '±0.6 ms'],
        ['GPS+PPS (ground truth)', '−0.001 ms', '±0.094 ms'],
    ])

    add_bullet_list(slide, 6.8, 4.1, 5.8, 3.0, [
        "HF timing is 100× better than internet NTP",
        "L2 has 3.3× tighter bound than L1",
        "Fusion D_clock: median −1.2 ms, 86% within ±5 ms",
        "A Tier 2 station bootstraps itself to Tier 3",
    ], font_size=15, color=(0x1b, 0x5e, 0x20))

    # D_clock histogram
    if os.path.exists(fig('fig3_fusion_dclock_histogram.png')):
        slide.shapes.add_picture(
            fig('fig3_fusion_dclock_histogram.png'),
            Inches(0.3), Inches(5.5), Inches(6.2), Inches(1.8))

    set_notes(slide,
        "Since we're listening to time standard stations and we know their "
        "broadcast schedules, we can recover UTC from the signals themselves. "
        "We run two independent Kalman filters — L1 uses geometric path "
        "delays, L2 adds an ionospheric correction. Both feed Chrony as "
        "reference clocks. The result: our HF-derived time is 100 times "
        "better than internet NTP and within about 1 millisecond of GPS "
        "ground truth. The ionospheric correction makes L2 three times "
        "tighter than L1 — it's doing real work. This is the bridge: a "
        "station with only a GPSDO can recover absolute time from the HF "
        "signals and unlock the D_clock products.")

    # =================================================================
    # SLIDE 5: Shared-Channel Discrimination
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "Shared-Channel Discrimination"

    add_textbox(slide, 0.4, 0.9, 6.0, 0.5,
                "Three stations on one frequency — physics tells them apart",
                font_size=16, bold=True, color=(0x15, 0x65, 0xc0))

    # Methods
    add_bullet_list(slide, 0.4, 1.5, 5.8, 2.5, [
        "Tick frequency: 1000 Hz (WWV/BPM) vs 1200 Hz (WWVH)",
        "Tick duration: 5 ms (WWV/WWVH) vs 10 ms (BPM)",
        "NIST tone schedule: ground truth 14 min/hr",
        "D_clock ordering: WWV < WWVH < BPM",
    ], font_size=14)

    # Evidence table
    add_table(slide, 0.4, 3.8, 5.8, 2.0, [
        ['Station', 'Records/day', 'Median D_clock', 'Median SNR'],
        ['WWV', '1,380', '−1.12 ms', '26.2 dB'],
        ['WWVH', '1,379', '−0.08 ms', '17.8 dB'],
        ['BPM', '1,380', '+1.48 ms', '21.1 dB'],
    ])

    add_textbox(slide, 0.4, 5.9, 5.8, 1.0,
                "Cross-station Doppler: r ≈ 0 (independent paths)\n"
                "Same-station cross-freq: r = 0.43 (shared path)",
                font_size=13, color=(0x44, 0x44, 0x44))

    # Discrimination figure
    if os.path.exists(fig('fig9_shared_channel_discrimination.png')):
        slide.shapes.add_picture(
            fig('fig9_shared_channel_discrimination.png'),
            Inches(6.5), Inches(0.9), Inches(6.5), Inches(6.2))

    set_notes(slide,
        "On the shared frequencies, three stations transmit simultaneously. "
        "We separate them with a layered approach. The strongest discriminator "
        "is the tick frequency gate — WWV at 1000 Hz, WWVH at 1200 Hz. For "
        "ground truth, the NIST tone schedule gives us 14 minutes per hour "
        "where only one station is broadcasting its audio tone. The physical "
        "validation is compelling: the D_clock offsets follow propagation "
        "geometry on all four shared frequencies, and the cross-station "
        "Doppler correlations are zero — proving these really are independent "
        "ionospheric paths, not artifacts of the discrimination.")

    # =================================================================
    # SLIDE 6: Carrier-Phase dTEC
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "Carrier-Phase dTEC — The Primary Science Product"

    add_textbox(slide, 0.4, 0.9, 6.0, 0.5,
                "Bypassing the propagation model noise floor",
                font_size=16, bold=True, color=(0x1b, 0x5e, 0x20))

    add_bullet_list(slide, 0.4, 1.5, 6.0, 3.0, [
        "Group-delay TEC: signal 0.85 ms, noise 6.5 ms → SNR 0.13 (buried)",
        "Carrier-phase dTEC: ~6 mTECU/min sensitivity → SNR 17–330×",
        "dTEC/dt = −f_D × c × f / 40.3",
        "",
        "This is a Tier 2 product — needs only the GPSDO",
        "",
        "17,045 dTEC records/day (per-minute)",
        "848,599 per-tick records/day (1-second resolution)",
        "GNSS-anchored: ZED-F9P VTEC 41.7 TECU, ±1 TECU",
    ], font_size=15)

    # dTEC figure
    if os.path.exists(fig('fig4_dtec_rate_timeseries.png')):
        slide.shapes.add_picture(
            fig('fig4_dtec_rate_timeseries.png'),
            Inches(6.5), Inches(0.9), Inches(6.5), Inches(6.2))

    set_notes(slide,
        "Here's the payoff. Group-delay TEC — the classical approach of "
        "measuring 1/f² dispersion — is below our noise floor. The "
        "propagation model has 6.5 millisecond errors, the dispersion "
        "signal is 0.85 milliseconds. We can't see it. But carrier-phase "
        "dTEC bypasses this entirely. We measure the Doppler shift — the "
        "rate of change of carrier phase — and convert it to dTEC/dt. The "
        "sensitivity is 6 milli-TECU per minute. And crucially, this is a "
        "Tier 2 product: it only needs the GPSDO, not absolute time.")

    # =================================================================
    # SLIDE 7: Differential dTEC
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "Differential dTEC — Self-Consistency"

    add_textbox(slide, 0.4, 0.9, 6.0, 0.5,
                "Same ionosphere, different frequencies — do they agree?",
                font_size=16, bold=True, color=(0x1b, 0x5e, 0x20))

    add_table(slide, 0.4, 1.8, 5.5, 1.5, [
        ['Station', 'Widest freq pair', 'RMS'],
        ['CHU', '3.33 – 14.67 MHz', '0.005–0.007 TECU'],
        ['WWV', '2.50 – 25.00 MHz', '0.005–0.026 TECU'],
    ])

    add_bullet_list(slide, 0.4, 3.5, 5.5, 2.0, [
        "22,474 records/day, all GOOD quality",
        "Also a Tier 2 product",
        "Cross-frequency consistency validates the physics",
    ], font_size=15, color=(0x1b, 0x5e, 0x20))

    # Differential dTEC figure
    if os.path.exists(fig('fig5_differential_dtec.png')):
        slide.shapes.add_picture(
            fig('fig5_differential_dtec.png'),
            Inches(6.2), Inches(0.9), Inches(6.8), Inches(6.2))

    set_notes(slide,
        "How do we know we're measuring real ionospheric physics? For each "
        "station, we compare dTEC at different frequencies on the same path. "
        "They agree to within 0.03 TECU RMS. This is 22,000 consistency "
        "checks per day, all passing.")

    # =================================================================
    # SLIDE 8: Physics Cascade (fig13)
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "Cross-Domain Consistency — Physics Cascade"

    # Full-width figure
    slide.shapes.add_picture(
        fig('fig13_physics_cascade.png'),
        Inches(0.3), Inches(0.8), Inches(7.5), Inches(6.4))

    # Right annotations
    add_textbox(slide, 8.0, 0.9, 5.0, 0.5,
                "CHU 7.85 MHz — exclusive channel",
                font_size=16, bold=True, color=(0x15, 0x65, 0xc0))

    add_bullet_list(slide, 8.0, 1.5, 5.0, 4.5, [
        "Panel A: D_clock (absolute delay, Tier 3+)",
        "Panel B: Doppler (phase rate, Tier 2)",
        "Panel C: dTEC/dt (derived from Doppler)",
        "Panel D: Integrated Doppler vs D_clock",
        "",
        "The money shot (Panel D):",
        "Integrated Doppler tracks the shape",
        "of D_clock at 82× smaller amplitude",
        "r = 0.60 correlation",
        "",
        "Same ionosphere, different physics,",
        "consistent measurements",
    ], font_size=14)

    set_notes(slide,
        "This is CHU at 7.85 MHz — an exclusive channel with no "
        "discrimination ambiguity. Four panels, one ionospheric path. "
        "At top, D_clock shows the diurnal propagation delay variation. "
        "Below, Doppler — the carrier phase slope. Then dTEC/dt derived "
        "from Doppler. The bottom panel is the key: if we integrate the "
        "Doppler, it should track the shape of D_clock. It does — r = 0.60 "
        "— but at 82 times smaller amplitude. The Doppler is measuring real "
        "ionospheric motion, at a sensitivity 82 times finer than D_clock "
        "noise.")

    # =================================================================
    # SLIDE 9: 17 Paths
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "17 Simultaneous Ionospheric Paths"

    # Ionospheric fingerprint figure (left)
    slide.shapes.add_picture(
        fig('fig10_ionospheric_fingerprint.png'),
        Inches(0.2), Inches(0.8), Inches(6.5), Inches(3.3))

    # Correlation heatmap (left bottom)
    slide.shapes.add_picture(
        fig('fig12_correlation_heatmap.png'),
        Inches(0.2), Inches(4.2), Inches(6.5), Inches(3.0))

    # Frequency ladder (right)
    slide.shapes.add_picture(
        fig('fig14_frequency_ladder.png'),
        Inches(6.8), Inches(0.8), Inches(6.3), Inches(3.3))

    # Doppler scatter triptych (right bottom)
    slide.shapes.add_picture(
        fig('fig11_doppler_scatter_triptych.png'),
        Inches(6.8), Inches(4.2), Inches(6.3), Inches(3.0))

    set_notes(slide,
        "The full system monitors 17 simultaneous paths through the "
        "ionosphere. On the shared channels, three stations at the same "
        "frequency give three independent ionospheric soundings — the "
        "Doppler correlation is zero, confirming they see different paths. "
        "Across frequencies, WWV from 2.5 to 25 MHz samples six different "
        "layers. CHU across three frequencies shows correlated Doppler — "
        "same path, same ionosphere, consistent. The correlation heatmap "
        "makes the structure clear: station clusters are independent, "
        "within-station cross-frequency is correlated. This is essentially "
        "a passive oblique ionosonde with 17 beams.")

    # =================================================================
    # SLIDE 10: What's Next
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "What's Next — Honest Assessment"

    # Left column: limits + deployed
    add_textbox(slide, 0.4, 0.9, 6.0, 0.4,
                "Current limits:", font_size=16, bold=True,
                color=(0xdc, 0x35, 0x45))
    add_bullet_list(slide, 0.4, 1.4, 6.0, 1.5, [
        "Group-delay TEC: below noise floor (SNR 0.13)",
        "VTEC maps from HF alone: noise-dominated",
        "Scintillation: infrastructure ready, awaiting storm",
    ], font_size=14)

    add_textbox(slide, 0.4, 2.8, 6.0, 0.4,
                "Just deployed:", font_size=16, bold=True,
                color=(0x1b, 0x5e, 0x20))
    add_bullet_list(slide, 0.4, 3.2, 6.0, 0.8, [
        "✅ GNSS-anchored dTEC (ZED-F9P, ±1 TECU)",
    ], font_size=14)

    add_textbox(slide, 0.4, 3.9, 6.0, 0.4,
                "Under development:", font_size=16, bold=True,
                color=(0x15, 0x65, 0xc0))
    add_bullet_list(slide, 0.4, 4.3, 6.0, 0.8, [
        "Tier 4: PPS injection into HF IQ stream",
        "Per-path slant correction for GNSS anchoring",
    ], font_size=14)

    # Right column: speculation
    add_textbox(slide, 6.8, 0.9, 6.0, 0.4,
                "Network of stations →", font_size=16, bold=True,
                color=(0xFF, 0x8C, 0x00))
    add_bullet_list(slide, 6.8, 1.4, 5.8, 2.0, [
        "Spatial TEC gradients → horizontal structure",
        "TID wavefront tracking (direction + velocity)",
        "Continental-scale passive oblique ionosonde",
        "No transmitter needed — use existing infrastructure",
    ], font_size=14)

    add_textbox(slide, 6.8, 3.6, 6.0, 0.4,
                "2–4 GPSDO-locked RX888s at one site →", font_size=16,
                bold=True, color=(0xFF, 0x8C, 0x00))
    add_bullet_list(slide, 6.8, 4.1, 5.8, 3.0, [
        "Phased-array angle-of-arrival (λ/2 ≈ 15 m at 10 MHz)",
        "Separate multipath by direction, not just delay",
        "Per-mode Doppler and dTEC on individual ray paths",
        "Scintillation spatial coherence measurements",
        "Same GPSDO → zero relative timing error",
    ], font_size=14)

    set_notes(slide,
        "What doesn't work yet: group-delay TEC is buried in noise. VTEC "
        "maps from HF alone aren't credible. Scintillation monitoring is "
        "built but the ionosphere has been quiet. What just shipped: "
        "GNSS-anchored dTEC, using a local GPS receiver to provide absolute "
        "scale. What's coming: PPS injection directly into the HF IQ stream "
        "— that's Tier 4, which would give us microsecond timing on every "
        "sample and potentially rescue group-delay TEC. But I want to leave "
        "you with two bigger ideas. First: a network of these stations. Each "
        "one gives 17 ionospheric paths. Ten stations across the continent "
        "gives 170 paths — that's a passive oblique ionosonde network using "
        "transmitters that are already on the air. Correlated Doppler across "
        "sites gives you TID wavefront direction and velocity. Second: "
        "multiple GPSDO-locked RX888s at a single site. Because they share "
        "the same 10 MHz reference, the antennas are phase-coherent — you "
        "get a phased array for free. Antenna spacing of 15 meters at 10 "
        "MHz gives you angle-of-arrival discrimination. That means you can "
        "separate multipath arrivals by direction, not just by delay, and "
        "track Doppler and dTEC on individual ray paths. The infrastructure "
        "is the same — the GPSDO is doing the heavy lifting.")

    # =================================================================
    # SLIDE 11: Summary & Call to Action
    # =================================================================
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text = "Summary & Call to Action"

    add_textbox(slide, 0.4, 1.0, 12.0, 0.6,
                "With an RX888 (~$180) and a GPSDO (~$162), you can:",
                font_size=22, bold=True, color=(0x1b, 0x5e, 0x20))

    add_bullet_list(slide, 0.6, 1.8, 11.5, 2.5, [
        "Measure ionospheric Doppler at 99.7% coverage, 24/7",
        "Extract dTEC/dt at ~6 mTECU/min sensitivity on 17 paths",
        "Discriminate three co-channel stations via physics",
        "Self-recover UTC to ±1 ms from the time signals",
    ], font_size=18, color=(0x22, 0x22, 0x22))

    add_textbox(slide, 0.6, 4.0, 11.5, 0.5,
                "Daily output: 24K timing · 850K phase · 17K dTEC · "
                "22K consistency checks",
                font_size=16, color=(0x44, 0x44, 0x44))

    add_textbox(slide, 0.6, 4.8, 11.5, 0.5,
                "Open source: github.com/mijahauan/hf-timestd  (MIT license)",
                font_size=18, bold=True, color=(0x21, 0x76, 0xFF))

    add_textbox(slide, 0.6, 5.8, 11.5, 1.0,
                "~$340 of hardware, open-source software, and the time "
                "standard stations already on the air → a 17-path "
                "ionospheric sounder running 24/7",
                font_size=20, bold=True, color=(0x1b, 0x5e, 0x20),
                alignment=PP_ALIGN.CENTER)

    set_notes(slide,
        "So: with about 340 dollars of hardware — an RX888 and a Leo Bodnar "
        "GPSDO — open-source software, and the time standard stations that "
        "are already on the air, you can build a 17-path ionospheric sounder "
        "that runs 24/7. The data products are scientifically meaningful, "
        "self-consistent, and validated against GPS ground truth. The code "
        "is on GitHub under MIT license. I'd love to see a network of these "
        "stations.")

    # =================================================================
    # Save
    # =================================================================
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    build_presentation()
